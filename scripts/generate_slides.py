# filename: scripts/generate_slides.py
# purpose:  Generate project presentation slides (.pptx)
# version:  1.0

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette ──────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1B, 0x3A, 0x5C)   # slide headers / title bg
TEAL    = RGBColor(0x17, 0x7E, 0x89)   # accent / section labels
GREEN   = RGBColor(0x27, 0xAE, 0x60)   # ✅ highlights
ORANGE  = RGBColor(0xE6, 0x7E, 0x22)   # ⚠️ / callout
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF4, 0xF6, 0xF8)   # slide background
DARK    = RGBColor(0x2C, 0x3E, 0x50)   # body text
GRAY    = RGBColor(0x7F, 0x8C, 0x8D)   # secondary text

W = Inches(13.33)   # 16:9 widescreen width
H = Inches(7.5)     # 16:9 widescreen height


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color: RGBColor = LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color: RGBColor):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txb(slide, text, l, t, w, h,
        size=18, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return box


def header_bar(slide, title, subtitle=None):
    """Navy bar across top with white title text."""
    add_rect(slide, 0, 0, W, Inches(1.35), NAVY)
    txb(slide, title,
        Inches(0.4), Inches(0.1), Inches(12), Inches(0.85),
        size=28, bold=True, color=WHITE)
    if subtitle:
        txb(slide, subtitle,
            Inches(0.4), Inches(0.9), Inches(12), Inches(0.4),
            size=14, color=RGBColor(0xAD, 0xD8, 0xE6))


def bullet_box(slide, items, l, t, w, h,
               size=17, color=DARK, indent=False):
    """Render a list of strings as bullet points."""
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = ("    " if indent else "") + item
        run.font.size  = Pt(size)
        run.font.color.rgb = color


def tag(slide, label, l, t, color=TEAL):
    """Small coloured label pill."""
    add_rect(slide, l, t, Inches(2.1), Inches(0.32), color)
    txb(slide, label,
        l + Inches(0.08), t + Inches(0.02), Inches(2.0), Inches(0.3),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    s = blank_slide(prs)
    bg(s, NAVY)
    # White gradient strip across middle
    add_rect(s, 0, Inches(2.1), W, Inches(3.3), RGBColor(0x0D, 0x25, 0x40))

    txb(s, "AI-Powered Customer Support Intelligence Platform",
        Inches(0.6), Inches(2.2), Inches(12), Inches(1.2),
        size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txb(s, "End-to-End ML System  |  NLP · Tabular ML · MLOps · Monitoring · Dashboard",
        Inches(0.6), Inches(3.45), Inches(12), Inches(0.5),
        size=16, color=RGBColor(0xAD, 0xD8, 0xE6), align=PP_ALIGN.CENTER)

    txb(s, "Python · DistilBERT · XGBoost · LightGBM · FastAPI · Airflow · Docker · Plotly Dash",
        Inches(0.6), Inches(3.95), Inches(12), Inches(0.45),
        size=13, color=RGBColor(0x85, 0xC1, 0xE9), align=PP_ALIGN.CENTER)

    txb(s, "PriyaMonisha  |  GUVI Data Science  |  2026",
        Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
        size=13, color=GRAY, align=PP_ALIGN.CENTER)


def slide_agenda(prs):
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "Agenda")

    cols = [
        ["1. Problem Statement", "2. Dataset Overview", "3. System Architecture",
         "4. EDA & Text Preprocessing", "5. Feature Engineering & Split"],
        ["6. ML Models & Results", "7. DistilBERT Fine-Tuning", "8. Regression & Clustering",
         "9. Explainability (SHAP + Attention)", "10. MLflow & Model Registry"],
        ["11. FastAPI Serving Layer", "12. Monitoring Stack", "13. Airflow Orchestration",
         "14. Docker Compose", "15. Plotly Dash Dashboard",
         "16. Key Findings & Conclusion"],
    ]
    xs = [Inches(0.4), Inches(4.6), Inches(8.8)]
    for col, x in zip(cols, xs):
        bullet_box(s, ["• " + c for c in col],
                   x, Inches(1.5), Inches(4.0), Inches(5.5), size=16)


def slide_problem(prs):
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "Problem Statement",
               "Manual ticket triage is slow, inconsistent, and expensive to scale")

    add_rect(s, Inches(0.4), Inches(1.5), Inches(3.9), Inches(5.5), WHITE)
    add_rect(s, Inches(4.5), Inches(1.5), Inches(3.9), Inches(5.5), WHITE)
    add_rect(s, Inches(8.6), Inches(1.5), Inches(4.3), Inches(5.5), WHITE)

    for x, title, color, items in [
        (Inches(0.4), "Task 1 — NLP", TEAL,
         ["Ticket Type Classification",
          "5 classes: Billing / Technical /",
          "Cancellation / Product / Refund",
          "Model: DistilBERT fine-tuned",
          "Input: Subject + Description"]),
        (Inches(4.5), "Task 2 — Tabular ML", NAVY,
         ["Ticket Priority Prediction",
          "4 classes: Low / Medium /",
          "High / Critical",
          "Model: XGBoost + Optuna",
          "Input: 17 tabular features"]),
        (Inches(8.6), "Task 3 — Regression", GREEN,
         ["Resolution Time Estimation",
          "Predict hours-to-resolve",
          "Model: RF / XGB / LGBM",
          "Input: tabular + text features",
          "Metric: RMSE, MAE, R²"]),
    ]:
        add_rect(s, x, Inches(1.5), Inches(3.8), Inches(0.42), color)
        txb(s, title, x + Inches(0.1), Inches(1.52),
            Inches(3.6), Inches(0.4), size=14, bold=True, color=WHITE)
        bullet_box(s, items, x + Inches(0.15), Inches(2.0),
                   Inches(3.6), Inches(4.8), size=14)


def slide_dataset(prs):
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "Dataset Overview",
               "Kaggle — Customer Support Ticket Dataset · 8,469 rows × 17 columns")

    stats = [
        ("8,469", "Total tickets"),
        ("17", "Raw columns"),
        ("5", "Ticket Type classes"),
        ("4", "Priority classes"),
        ("70 / 10 / 20", "Train / Val / Test split"),
        ("1,931", "Closed tickets (regression subset)"),
    ]
    for i, (val, label) in enumerate(stats):
        col = i % 3
        row = i // 3
        x = Inches(0.5 + col * 4.3)
        y = Inches(1.65 + row * 2.2)
        add_rect(s, x, y, Inches(3.9), Inches(1.9), WHITE)
        txb(s, val, x + Inches(0.1), y + Inches(0.15),
            Inches(3.7), Inches(1.0), size=34, bold=True, color=NAVY,
            align=PP_ALIGN.CENTER)
        txb(s, label, x + Inches(0.1), y + Inches(1.05),
            Inches(3.7), Inches(0.6), size=14, color=GRAY,
            align=PP_ALIGN.CENTER)

    bullet_box(s,
        ["⚠  5 Ticket Type classes found — spec said 4 (data-first validation caught this)",
         "⚠  5,700 structural nulls — all tied to Ticket Status ≠ 'Closed' (not data quality issues)",
         "⚠  {product_purchased} placeholder in all descriptions — replaced during cleaning",
         "⚠  1,365 negative hours_to_resolve — fixed with +24h timestamp wrap-around correction"],
        Inches(0.4), Inches(6.1), Inches(12.5), Inches(1.2), size=13, color=ORANGE)


def slide_architecture(prs):
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "System Architecture — End-to-End ML Pipeline")

    layers = [
        ("Data Layer",       "Raw CSV → clean.py (8-fn pipeline) → Pandera validation → PostgreSQL + Redis cache",         TEAL),
        ("Feature Layer",    "TextPreprocessor (VADER + TF-IDF) · TabularEncoder (Ordinal + TargetEncoder) · 17 features",  NAVY),
        ("Model Layer",      "Baselines (LR/NB/DT) → Advanced (RF/XGB/LGBM + Optuna) → DistilBERT (Colab T4)",             TEAL),
        ("MLOps Layer",      "MLflow tracking (47 runs) · Model Registry (Production/Challenger aliases) · Optuna HPO",     NAVY),
        ("Serving Layer",    "FastAPI — 8 endpoints · X-API-Key auth · /admin/reload hot-swap · Prometheus metrics",        TEAL),
        ("Monitoring Layer", "Prometheus + Grafana + Alertmanager · Evidently PSI drift · 5 alert rules · RUNBOOK.md",     NAVY),
        ("Orchestration",    "Airflow — 4 DAGs (ETL, drift monitor, retrain, model report) · 26 tasks total",              TEAL),
        ("Deployment",       "Docker Compose — 10 containers on one bridge network (csip-net) · CI: GitHub Actions",       NAVY),
        ("Dashboard",        "Plotly Dash — 6 pages · EDA, Leaderboard, Live Predictions, Drift, Clustering/SHAP",         TEAL),
    ]
    for i, (layer, desc, color) in enumerate(layers):
        y = Inches(1.45 + i * 0.665)
        add_rect(s, Inches(0.3), y, Inches(2.1), Inches(0.56), color)
        txb(s, layer, Inches(0.35), y + Inches(0.08),
            Inches(2.0), Inches(0.45), size=12, bold=True, color=WHITE)
        add_rect(s, Inches(2.5), y, Inches(10.4), Inches(0.56), WHITE)
        txb(s, desc, Inches(2.6), y + Inches(0.1),
            Inches(10.2), Inches(0.45), size=12, color=DARK)


def slide_eda(prs):
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "EDA & Text Preprocessing",
               "14 charts · VADER sentiment · TF-IDF · spaCy lemmatization")

    add_rect(s, Inches(0.4), Inches(1.5), Inches(6.0), Inches(2.7), WHITE)
    txb(s, "Key EDA Findings", Inches(0.5), Inches(1.55),
        Inches(5.8), Inches(0.4), size=15, bold=True, color=NAVY)
    bullet_box(s, [
        "• Ticket Type: 5 classes, nearly balanced (1.07:1 ratio)",
        "• Ticket Priority: 4 classes, ~2,000 tickets each",
        "• 80% of Critical tickets breach their 4-hour SLA window",
        "• CSAT: mean 2.99/5, but 67.3% null (open tickets have no rating)",
        "• All numeric CSAT correlations < 0.05 — no exploitable signal",
    ], Inches(0.5), Inches(1.95), Inches(5.8), Inches(2.1), size=14)

    add_rect(s, Inches(6.7), Inches(1.5), Inches(6.2), Inches(2.7), WHITE)
    txb(s, "Text Preprocessing Pipeline", Inches(6.8), Inches(1.55),
        Inches(6.0), Inches(0.4), size=15, bold=True, color=NAVY)
    bullet_box(s, [
        "• Lowercase, remove punctuation, HTML artifacts, {placeholder} strings",
        "• Tokenize → stopword removal → spaCy lemmatization",
        "• VADER sentiment: 72.7% positive / 4.3% neutral / 23.1% negative",
        "• TF-IDF: 9,307-term vocab (min_df=2 on 5,928 training docs)",
        "• Top TF-IDF terms identical across all 5 Ticket Type classes ← key signal",
    ], Inches(6.8), Inches(1.95), Inches(6.0), Inches(2.1), size=14)

    add_rect(s, Inches(0.4), Inches(4.4), Inches(12.5), Inches(2.7), WHITE)
    txb(s, "Critical SLA Finding  →  The Business Case for This Platform",
        Inches(0.5), Inches(4.45), Inches(12), Inches(0.45),
        size=15, bold=True, color=ORANGE)
    bullet_box(s, [
        "• Resolution time is uniformly distributed 0–24 hours in this dataset",
        "• Critical priority SLA = 4 hours  →  P(breach) = 4/24 expected under random assignment",
        "• Actual breach rate = 80%  →  confirms Critical tickets are not being prioritised fast enough",
        "• Automated priority prediction (Task 2) directly addresses this gap — even modest accuracy lifts have material SLA impact",
    ], Inches(0.5), Inches(4.9), Inches(12.3), Inches(2.0), size=14)


def slide_features(prs):
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "Feature Engineering, Split & Redis Feature Store")

    add_rect(s, Inches(0.4), Inches(1.5), Inches(5.8), Inches(3.5), WHITE)
    txb(s, "TabularEncoder  (17 model-ready features)", Inches(0.5), Inches(1.55),
        Inches(5.6), Inches(0.4), size=15, bold=True, color=NAVY)
    bullet_box(s, [
        "• OrdinalEncoder  →  channel (4), gender (3)",
        "• TargetEncoder (multiclass)  →  product_purchased (5 cols, one per class)",
        "• Numeric passthrough  →  age, days_since_purchase, word_count,",
        "    char_count, sentiment score, response_hour_of_day, + 4 more",
        "• 4 constant columns auto-dropped (zero variance on training data)",
        "• Regression subset uses keep_mask  →  13 features (drops 4 text-only cols)",
    ], Inches(0.5), Inches(1.95), Inches(5.6), Inches(2.9), size=14)

    add_rect(s, Inches(6.5), Inches(1.5), Inches(6.4), Inches(3.5), WHITE)
    txb(s, "Train / Val / Test Split", Inches(6.6), Inches(1.55),
        Inches(6.2), Inches(0.4), size=15, bold=True, color=NAVY)
    bullet_box(s, [
        "• Stratified split preserving class proportions across all 3 sets",
        "• Classification:  5,928 / 847 / 1,694  (70 / 10 / 20 %)",
        "• Regression subset (closed tickets):  1,931 / 276 / 562",
        "",
        "Redis Feature Store:",
        "• SHA-256 key hashing, connection pool, pipeline writes",
        "• Graceful degrade — except redis.RedisError: return None",
        "• System runs correctly with or without Redis available",
    ], Inches(6.6), Inches(1.95), Inches(6.2), Inches(2.9), size=14)

    add_rect(s, Inches(0.4), Inches(5.2), Inches(12.5), Inches(2.0), WHITE)
    txb(s, "Why TargetEncoder for Product Purchased?",
        Inches(0.5), Inches(5.25), Inches(12), Inches(0.4),
        size=14, bold=True, color=TEAL)
    bullet_box(s, [
        "• Product Purchased has high cardinality (many unique values) — OrdinalEncoder would produce arbitrary integer codes with no semantic meaning",
        "• TargetEncoder replaces each category with the mean of each target class, creating 5 columns (one per Ticket Type class) with real predictive signal",
        "• Multiclass TargetEncoder requires get_feature_names_out() to name the output columns — a non-obvious sklearn API detail discovered during implementation",
    ], Inches(0.5), Inches(5.65), Inches(12.3), Inches(1.4), size=13)


def slide_models(prs):
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "ML Models & Results",
               "Baselines → Advanced Ensemble → Optuna Hyperparameter Tuning")

    # results table
    rows = [
        ("Algorithm", "Task", "Val F1-macro", "Notes", True),
        ("Logistic Regression (baseline)", "Ticket Type", "0.1913", "Near noise floor (5-class ≈ 0.20)", False),
        ("Naive Bayes (baseline)", "Ticket Priority", "0.2418", "Best baseline for priority", False),
        ("Random Forest + Optuna", "Ticket Type", "0.1795", "Below LR baseline", False),
        ("XGBoost + Optuna", "Ticket Priority", "0.2625 ✅", "Best overall — real signal above 0.25 floor", False),
        ("LightGBM + Optuna", "Ticket Type", "0.1997", "Best for Type — still at noise floor", False),
        ("Dummy (most-frequent)", "Ticket Type", "0.069", "True irreducible floor — all models beat this", False),
        ("Dummy (most-frequent)", "Ticket Priority", "0.103", "XGBoost 2.5× better than dummy", False),
    ]
    col_w = [Inches(3.3), Inches(2.2), Inches(2.0), Inches(4.8)]
    col_x = [Inches(0.3), Inches(3.7), Inches(6.0), Inches(8.1)]
    y = Inches(1.5)
    for row in rows:
        header_row = row[4]
        fill = NAVY if header_row else WHITE
        text_color = WHITE if header_row else DARK
        row_h = Inches(0.5) if header_row else Inches(0.62)
        for cell, cw, cx in zip(row[:4], col_w, col_x):
            add_rect(s, cx, y, cw, row_h, fill)
            txb(s, cell, cx + Inches(0.06), y + Inches(0.08),
                cw - Inches(0.1), row_h - Inches(0.1),
                size=13, bold=header_row, color=text_color)
        y += row_h

    txb(s, "6 Optuna studies  ·  MedianPruner  ·  class_weight='balanced'  ·  refit on full training data after tuning",
        Inches(0.3), Inches(6.8), Inches(12.7), Inches(0.4),
        size=13, color=GRAY)


def slide_distilbert(prs):
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "DistilBERT Fine-Tuning",
               "Google Colab T4  ·  5 epochs  ·  AdamW + linear LR warmup  ·  pair tokenization")

    add_rect(s, Inches(0.4), Inches(1.5), Inches(5.8), Inches(5.6), WHITE)
    txb(s, "Training Setup", Inches(0.5), Inches(1.55),
        Inches(5.6), Inches(0.4), size=15, bold=True, color=NAVY)
    bullet_box(s, [
        "• Input: Ticket Subject + [SEP] + Ticket Description",
        "• MAX_LENGTH = 128 tokens  ·  pre-tokenized once per split",
        "• 4 AdamW parameter groups (backbone vs head, bias/LayerNorm excluded from weight decay)",
        "• Differential LR: backbone lr=1e-5, head lr=1e-4",
        "• Early stopping patience=2  ·  atomic checkpointing (tempfile + shutil.move)",
        "• Saved with model.save_pretrained() to models/distilbert/",
    ], Inches(0.5), Inches(1.95), Inches(5.6), Inches(2.3), size=14)

    txb(s, "Test Set Results", Inches(0.5), Inches(4.35),
        Inches(5.6), Inches(0.4), size=15, bold=True, color=NAVY)
    for label, val in [
        ("F1-macro", "0.1954"),
        ("F1-weighted", "0.1960"),
        ("Accuracy", "0.2048"),
        ("Best val F1 (epoch 4)", "0.1823"),
    ]:
        y_off = Inches(4.75) + [label].index(label) * 0
    y = Inches(4.75)
    for label, val in [
        ("F1-macro", "0.1954"), ("F1-weighted", "0.1960"),
        ("Accuracy", "0.2048"), ("Best val F1 @ epoch 4", "0.1823"),
    ]:
        add_rect(s, Inches(0.5), y, Inches(3.5), Inches(0.42), LIGHT)
        add_rect(s, Inches(4.1), y, Inches(1.9), Inches(0.42), NAVY)
        txb(s, label, Inches(0.6), y + Inches(0.07), Inches(3.3), Inches(0.35), size=14)
        txb(s, val, Inches(4.15), y + Inches(0.07), Inches(1.7), Inches(0.35),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        y += Inches(0.47)

    add_rect(s, Inches(6.5), Inches(1.5), Inches(6.4), Inches(5.6), WHITE)
    txb(s, "Key Finding — The Noise Floor", Inches(6.6), Inches(1.55),
        Inches(6.2), Inches(0.4), size=15, bold=True, color=ORANGE)
    bullet_box(s, [
        "DistilBERT (0.1954) ≈ LightGBM (0.1997)",
        "≈ random guess floor (0.20 for 5 classes)",
        "",
        "val_loss RISES every epoch while train_loss barely falls",
        "→ classic 'no learnable signal' training curve",
        "",
        "Starting cross-entropy ≈ ln(5) = 1.609",
        "→ confirms correct wiring; model not stuck",
        "",
        "4 independent confirmations:",
        "1. TF-IDF top terms identical across all classes (§5.3)",
        "2. Classical ML baselines at noise floor (§5.5)",
        "3. Optuna-tuned ensembles at noise floor (§5.6)",
        "4. DistilBERT fine-tune confirms ceiling (§5.9)",
        "",
        "Most likely cause: Ticket Type labels are synthetic /",
        "near-randomly assigned in the source dataset.",
        "",
        "API response: model_status: 'below_quality_bar'",
        "+ reliability_note field on every /predict/type call",
    ], Inches(6.6), Inches(1.95), Inches(6.2), Inches(5.0), size=13)


def slide_regression_clustering(prs):
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "Regression (Resolution Time) & Unsupervised Clustering")

    add_rect(s, Inches(0.4), Inches(1.5), Inches(6.0), Inches(5.6), WHITE)
    txb(s, "Resolution Time Regression (Task 3)", Inches(0.5), Inches(1.55),
        Inches(5.8), Inches(0.4), size=15, bold=True, color=NAVY)
    bullet_box(s, [
        "Models: RF / XGBoost / LightGBM + Optuna (KFold CV, minimize RMSE)",
        "Optional log-transform: skewness-gated (skew=0.031 → not applied)",
        "Metrics: RMSE / MAE / R² / MAPE / RMSLE",
        "",
        "Results:",
    ], Inches(0.5), Inches(1.95), Inches(5.8), Inches(1.6), size=14)
    for m, rmse, r2 in [
        ("Random Forest", "7.09 hrs", "−0.01 to −0.05"),
        ("XGBoost",       "7.09 hrs", "−0.01 to −0.05"),
        ("LightGBM",      "7.11 hrs", "−0.01 to −0.05"),
        ("DummyRegressor","7.06 hrs", "baseline (mean predictor)"),
    ]:
        y_pos = Inches(3.6)
    y = Inches(3.6)
    for m, rmse, r2 in [
        ("Random Forest", "7.09", "R² ≈ −0.02"),
        ("XGBoost",       "7.09", "R² ≈ −0.01"),
        ("LightGBM",      "7.11", "R² ≈ −0.05"),
        ("Dummy mean",    "7.06", "beats all — confirms no signal"),
    ]:
        add_rect(s, Inches(0.5), y, Inches(2.8), Inches(0.42), LIGHT)
        add_rect(s, Inches(3.4), y, Inches(1.0), Inches(0.42),
                 ORANGE if "Dummy" in m else NAVY)
        add_rect(s, Inches(4.5), y, Inches(1.8), Inches(0.42), LIGHT)
        txb(s, m,    Inches(0.55), y + Inches(0.07), Inches(2.7), Inches(0.35), size=12)
        txb(s, rmse, Inches(3.45), y + Inches(0.07), Inches(0.9), Inches(0.35),
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txb(s, r2,   Inches(4.55), y + Inches(0.07), Inches(1.7), Inches(0.35), size=11)
        y += Inches(0.46)

    bullet_box(s, ["Conclusion: near-uniform 0–24h distribution carries no exploitable tabular signal."],
               Inches(0.5), Inches(5.65), Inches(5.8), Inches(0.8), size=13, color=ORANGE)

    add_rect(s, Inches(6.7), Inches(1.5), Inches(6.2), Inches(5.6), WHITE)
    txb(s, "K-Means Customer Segmentation", Inches(6.8), Inches(1.55),
        Inches(6.0), Inches(0.4), size=15, bold=True, color=NAVY)
    bullet_box(s, [
        "Swept K = 2 to 6  ·  StandardScaler  ·  n_init='auto'",
        "Best K = 2  (Silhouette = 0.1573)",
        "Davies-Bouldin Index = 2.1902  (also minimum at K=2)",
        "→ Both metrics agree: K=2 is the natural cluster count",
        "",
        "Visualization: PCA (2D projection) + t-SNE (non-linear)",
        "Stratified t-SNE subsample for speed",
        "",
        "All results logged to MLflow (csip-clustering experiment)",
        "Charts in Dash Page 5 — Clustering & Explainability",
        "",
        "Fairness analysis (Phase C):",
        "• XGBoost priority model evaluated by Gender / Age / Channel",
        "• All segment F1s within ±0.03 of overall (0.251)",
        "• Consistent with statistical noise — no systematic bias",
        "• Heatmaps in Dash fairness gallery",
    ], Inches(6.8), Inches(1.95), Inches(6.0), Inches(4.9), size=13)


def slide_shap(prs):
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "Model Explainability — SHAP + DistilBERT Attention")

    add_rect(s, Inches(0.4), Inches(1.5), Inches(6.0), Inches(5.6), WHITE)
    txb(s, "SHAP — TreeExplainer", Inches(0.5), Inches(1.55),
        Inches(5.8), Inches(0.4), size=15, bold=True, color=NAVY)
    bullet_box(s, [
        "Applied to: XGBoost Priority + RF Regression models",
        "Handles (N, F, C) 3D array — LGBM multi-class SHAP output",
        "Fixed: must index [:, :, class_idx] to get per-class attributions",
        "",
        "Charts generated:",
        "  • Beeswarm plot — global feature importance across all samples",
        "  • Waterfall plot — per-prediction breakdown",
        "  • Mid-confidence correct + high-confidence wrong examples selected",
        "",
        "Top driver — Priority prediction:",
        "  → days_since_purchase  (newer purchases → higher urgency)",
        "",
        "Top driver — Resolution time:",
        "  → response_hour_of_day  (time-of-day operational patterns)",
        "",
        "SHAP values logged to MLflow: csip-explainability experiment",
    ], Inches(0.5), Inches(1.95), Inches(5.8), Inches(5.0), size=13)

    add_rect(s, Inches(6.7), Inches(1.5), Inches(6.2), Inches(5.6), WHITE)
    txb(s, "DistilBERT Attention Heatmaps", Inches(6.8), Inches(1.55),
        Inches(6.0), Inches(0.4), size=15, bold=True, color=NAVY)
    bullet_box(s, [
        "Notebook: 09b_distilbert_attention.py",
        "Loads fine-tuned model with output_attentions=True",
        "300-sample subset (FAST_MODE), sample accuracy = 0.2433",
        "",
        "Method:",
        "  • Extract per-layer [CLS]-token attention",
        "  • Average over all attention heads",
        "  • Plot top-16 most-attended tokens as seaborn heatmap",
        "  • Separate heatmaps: highest-confidence correct vs wrong",
        "",
        "Finding:",
        "  • Correct predictions: 44% top tokens = special/punctuation",
        "  • Wrong predictions: 50% top tokens = special/punctuation",
        "  • Model attends to [CLS], [SEP], commas — not content words",
        "  → Mechanistic confirmation: DistilBERT cannot find signal",
        "     because the ticket text contains no learnable type cues",
        "",
        "Charts: artifacts/charts/distilbert_attention_*.png",
        "Displayed in Dash Page 5 gallery",
    ], Inches(6.8), Inches(1.95), Inches(6.0), Inches(4.9), size=13)


def slide_mlflow(prs):
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "MLflow Experiment Tracking & Model Registry")

    stats = [("47+", "Logged runs"), ("5", "Experiments"), ("3", "Registered models")]
    for i, (val, label) in enumerate(stats):
        x = Inches(0.4 + i * 4.3)
        add_rect(s, x, Inches(1.5), Inches(3.9), Inches(1.4), WHITE)
        txb(s, val, x + Inches(0.1), Inches(1.55), Inches(3.7), Inches(0.85),
            size=42, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txb(s, label, x + Inches(0.1), Inches(2.22), Inches(3.7), Inches(0.5),
            size=14, color=GRAY, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(0.4), Inches(3.1), Inches(12.5), Inches(4.0), WHITE)
    txb(s, "Model Registry  —  Production / Challenger Aliases",
        Inches(0.5), Inches(3.15), Inches(12), Inches(0.4),
        size=15, bold=True, color=NAVY)

    headers = ["Registered Model", "Production (v1)", "Metric", "Challenger (v2)"]
    rows_data = [
        ("csip-ticket-type-classifier", "LightGBM", "val F1-macro = 0.1997", "XGBoost"),
        ("csip-priority-classifier",    "XGBoost",  "val F1-macro = 0.2625", "Random Forest"),
        ("csip-resolution-regressor",   "Random Forest", "val RMSE = 7.07 hrs", "XGBoost"),
    ]
    col_w = [Inches(3.5), Inches(2.2), Inches(3.3), Inches(2.2)]
    col_x = [Inches(0.5), Inches(4.1), Inches(6.4), Inches(9.8)]
    y = Inches(3.6)
    for i, row in enumerate([headers] + rows_data):
        is_hdr = i == 0
        fill = TEAL if is_hdr else (LIGHT if i % 2 == 0 else WHITE)
        tc   = WHITE if is_hdr else DARK
        for cell, cw, cx in zip(row, col_w, col_x):
            add_rect(s, cx, y, cw, Inches(0.52), fill)
            txb(s, cell, cx + Inches(0.06), y + Inches(0.1),
                cw - Inches(0.1), Inches(0.4),
                size=13, bold=is_hdr, color=tc)
        y += Inches(0.52)

    bullet_box(s, [
        "Experiments: csip-baseline-classifiers · csip-advanced-ml · csip-regression-models · csip-clustering · csip-explainability",
        "Rules: mlflow.search_runs(experiment_ids=[...]) — never experiment_names=. All params cast: int(run['params.k']), float(run['metrics.f1'])",
    ], Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.65), size=12, color=GRAY)


def slide_fastapi(prs):
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "FastAPI Serving Layer",
               "8 endpoints · X-API-Key auth · /admin/reload hot-swap · Prometheus metrics · live-tested")

    endpoints = [
        ("GET  /health",             "—",                 "Liveness / readiness  (200 healthy · 503 loading)"),
        ("GET  /metrics",            "—",                 "Prometheus scrape endpoint (text/plain)"),
        ("POST /predict/type",       "CSIP_API_KEY",      "Ticket type + model_status + reliability_note"),
        ("POST /predict/priority",   "CSIP_API_KEY",      "Priority (Low/Medium/High/Critical)"),
        ("POST /predict/resolution", "CSIP_API_KEY",      "Hours-to-resolve estimate (RMSE ~7 hrs)"),
        ("POST /explain/priority",   "CSIP_API_KEY",      "SHAP-based explanation of priority prediction"),
        ("POST /admin/reload",       "ADMIN_API_KEY",     "Hot-swap models — lock-guarded, 600s rate limit"),
        ("POST /admin/drift-check",  "ADMIN_API_KEY",     "On-demand PSI drift check across 17 features"),
    ]
    y = Inches(1.5)
    for ep, auth, desc in endpoints:
        color = GREEN if "admin" not in ep.lower() else ORANGE
        add_rect(s, Inches(0.3), y, Inches(3.5), Inches(0.48), color)
        add_rect(s, Inches(3.9), y, Inches(2.2), Inches(0.48), LIGHT)
        add_rect(s, Inches(6.2), y, Inches(6.7), Inches(0.48), WHITE)
        txb(s, ep,   Inches(0.38), y + Inches(0.08), Inches(3.3), Inches(0.36),
            size=12, bold=True, color=WHITE)
        txb(s, auth, Inches(3.95), y + Inches(0.08), Inches(2.1), Inches(0.36), size=11)
        txb(s, desc, Inches(6.28), y + Inches(0.08), Inches(6.5), Inches(0.36), size=12)
        y += Inches(0.51)

    bullet_box(s, [
        "Observed latency (live-tested against running server):  /predict/type ~54ms  ·  /predict/priority ~14ms  ·  /predict/resolution ~15ms  ·  /explain/priority ~27ms  →  all < 500ms target",
        "Two-tier keys: ADMIN_API_KEY (superset) + CSIP_API_KEY (read-only, used by Dash so it never holds admin credentials)  ·  secrets.compare_digest (constant-time comparison, no timing attacks)",
        "lifespan context manager loads all 9 artifacts via run_in_executor (non-blocking event loop)  ·  startup failure → app.state.ready=False + /health 503 instead of container crash",
    ], Inches(0.3), Inches(5.7), Inches(12.7), Inches(1.7), size=12, color=DARK)


def slide_monitoring(prs):
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "Monitoring Stack",
               "Prometheus + Grafana + Alertmanager + Evidently PSI Drift Detection")

    add_rect(s, Inches(0.4), Inches(1.5), Inches(4.0), Inches(5.6), WHITE)
    txb(s, "Prometheus Metrics (5)", Inches(0.5), Inches(1.55),
        Inches(3.8), Inches(0.4), size=14, bold=True, color=NAVY)
    bullet_box(s, [
        "csip_predictions_total (counter)",
        "csip_prediction_latency_seconds (histogram)",
        "csip_prediction_errors_total (counter)",
        "csip_prediction_confidence (histogram)",
        "csip_models_loaded (gauge)",
        "+ 3 drift gauges (PSI per feature,",
        "  max PSI, drift detected flag)",
        "",
        "5 Alert Rules:",
        "HighPredictionErrorRate  (warn)",
        "HighPredictionLatencyP95 (warn)",
        "FeatureDriftDetected     (warn)",
        "ModelsNotLoaded          (critical)",
        "DriftCheckStale          (info)",
        "",
        "RUNBOOK.md: diagnosis +",
        "remediation for each alert",
    ], Inches(0.5), Inches(1.95), Inches(3.8), Inches(5.0), size=13)

    add_rect(s, Inches(4.6), Inches(1.5), Inches(4.1), Inches(5.6), WHITE)
    txb(s, "PSI Drift Detection (Evidently)", Inches(4.7), Inches(1.55),
        Inches(3.9), Inches(0.4), size=14, bold=True, color=NAVY)
    bullet_box(s, [
        "src/monitoring/drift.py:",
        "  • load_baseline() — reads training dist",
        "  • _psi_categorical() / _psi_continuous()",
        "  • check_drift() — returns per-feature PSI",
        "",
        "Cross-validated vs Evidently 0.4.30:",
        "  • Negative control: 0/17 features drifted",
        "  • Positive control (+2σ on Customer Age):",
        "    1/17 drifted — both methods agree",
        "",
        "Threshold: PSI ≥ 0.10 = moderate drift",
        "Live test (held-out test split):",
        "  max_psi = 0.088 → no drift ✅",
        "",
        "16 pytest tests cover all PSI logic",
        "",
        "POST /admin/drift-check on demand",
        "Daily csip_drift_monitor Airflow DAG",
    ], Inches(4.7), Inches(1.95), Inches(3.9), Inches(5.0), size=13)

    add_rect(s, Inches(8.9), Inches(1.5), Inches(4.0), Inches(5.6), WHITE)
    txb(s, "Grafana Dashboard (6 panels)", Inches(9.0), Inches(1.55),
        Inches(3.8), Inches(0.4), size=14, bold=True, color=NAVY)
    bullet_box(s, [
        "Auto-provisioned via:",
        "monitoring/grafana/provisioning/",
        "  datasources/ + dashboards/",
        "",
        "Panels:",
        "  • Prediction rate by task",
        "  • P95 latency over time",
        "  • Error rate by task",
        "  • Confidence distribution",
        "  • Models loaded gauge",
        "  • Feature drift PSI heatmap",
        "",
        "Alertmanager routing:",
        "  critical → slack-notifications",
        "  (1h repeat interval)",
        "  warning/info → UI only",
        "",
        "Airflow has independent",
        "  CSIP_ALERT_WEBHOOK path",
    ], Inches(9.0), Inches(1.95), Inches(3.8), Inches(5.0), size=13)


def slide_airflow(prs):
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "Airflow Orchestration — 4 DAGs, 26 Tasks Total")

    dags = [
        ("csip_etl", "5 tasks · Daily 02:00 UTC",
         ["clean → validate → postgres_upsert", "→ regen_feature_arrays → done",
          "Writes all 8,469 rows to PostgreSQL tickets table",
          "NaT→None fix for nullable columns"], TEAL),
        ("csip_drift_monitor", "6 tasks · Daily 03:00 UTC",
         ["load_baseline → check_psi →", "2-branch: no_drift_log | alert_drift",
          "Reads PSI results from FastAPI /admin/drift-check",
          "Sends webhook alert if drift detected"], NAVY),
        ("csip_retrain", "10 tasks · Sunday 02:00 UTC · max_active_runs=1",
         ["fan-out: retrain LGBM / XGB / RF in parallel",
          "3 ShortCircuitOperator PSI gates (per task group)",
          "→ evaluate → 3-branch: promote+regen_baseline | skip | alert",
          "Atomic shutil.move promotion, .json.bak registry backup"], TEAL),
        ("csip_model_report", "5 tasks · Sunday 06:00 UTC",
         ["MLflow pull → build leaderboard →", "update registry JSON → write report",
          "Reads 47+ runs via search_runs(experiment_ids=[...])",
          "Generates weekly model comparison report"], NAVY),
    ]
    y = Inches(1.5)
    for name, schedule, items, color in dags:
        add_rect(s, Inches(0.3), y, Inches(2.8), Inches(1.6), color)
        txb(s, name, Inches(0.38), y + Inches(0.1),
            Inches(2.6), Inches(0.5), size=13, bold=True, color=WHITE)
        txb(s, schedule, Inches(0.38), y + Inches(0.58),
            Inches(2.6), Inches(0.9), size=11, color=RGBColor(0xAD, 0xD8, 0xE6))
        add_rect(s, Inches(3.2), y, Inches(9.7), Inches(1.6), WHITE)
        bullet_box(s, ["• " + it for it in items],
                   Inches(3.3), y + Inches(0.05), Inches(9.5), Inches(1.5), size=13)
        y += Inches(1.7)

    bullet_box(s, [
        "Key decisions: storage=None Optuna studies (no SQLite lock contention in parallel fan-out)  ·  float() XCom cast (np.float64 not JSON-safe)  ·  trigger_rule=NONE_FAILED_MIN_ONE_SUCCESS on evaluate task  ·  fail-open PSI gates (missing metrics file → always retrain)"
    ], Inches(0.3), Inches(7.15), Inches(12.7), Inches(0.35), size=12, color=GRAY)


def slide_docker(prs):
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "Docker Compose — 10-Container Full-Stack Deployment")

    services = [
        ("postgres",           "PostgreSQL 15",          "csip + airflow databases (init-multi-db.sh)",         TEAL),
        ("redis",              "Redis 7",                "Feature cache  ·  graceful degrade",                  GRAY),
        ("mlflow",             "MLflow 2.13.2",          "Experiment tracking at :5001  ·  custom Dockerfile",  TEAL),
        ("fastapi",            "FastAPI (root image)",   "Full requirements.txt  ·  libgomp1 for LightGBM",     NAVY),
        ("dash",               "Plotly Dash (gunicorn)", "Trimmed requirements  ·  CSIP_API_KEY (read-only)",   NAVY),
        ("prometheus",         "Prometheus",             "Scrapes :8000/metrics  ·  5 alert rules",             TEAL),
        ("alertmanager",       "Alertmanager",           "Slack webhook (critical)  ·  UI only (warning)",      GRAY),
        ("grafana",            "Grafana 10",             "6-panel dashboard  ·  auto-provisioned",              TEAL),
        ("airflow-webserver",  "Airflow 2.9.3",          "Constraint-based pip installs  ·  LocalExecutor",     NAVY),
        ("airflow-scheduler",  "Airflow Scheduler",      "All 4 DAGs  ·  selective bind mounts for src/",       NAVY),
    ]
    y = Inches(1.5)
    for i, (name, image, desc, color) in enumerate(services):
        row_h = Inches(0.565)
        fill = LIGHT if i % 2 == 0 else WHITE
        add_rect(s, Inches(0.3), y, Inches(2.3), row_h, color)
        add_rect(s, Inches(2.7), y, Inches(2.5), row_h, fill)
        add_rect(s, Inches(5.3), y, Inches(7.6), row_h, fill)
        txb(s, name,  Inches(0.38), y + Inches(0.1), Inches(2.1), row_h - Inches(0.1),
            size=12, bold=True, color=WHITE)
        txb(s, image, Inches(2.76), y + Inches(0.1), Inches(2.3), row_h - Inches(0.1), size=12)
        txb(s, desc,  Inches(5.36), y + Inches(0.1), Inches(7.4), row_h - Inches(0.1), size=12)
        y += row_h

    bullet_box(s, [
        "All services: restart: unless-stopped  ·  single bridge network csip-net  ·  model artifacts bind-mounted (not baked) so /admin/reload and Airflow promotions share the same files",
        "Real bugs caught during live bring-up: host port 5432 shadowed by native Windows PostgreSQL service → remapped 5433:5432  ·  postgres container recreate crashes airflow-scheduler (SQLAlchemy persistent conn) → manual restart needed",
    ], Inches(0.3), Inches(7.18), Inches(12.7), Inches(0.65), size=12, color=GRAY)


def slide_dashboard(prs):
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "Plotly Dash Dashboard — 6 Pages",
               "use_pages=True  ·  def layout() per page (re-reads artifacts on every navigation, zero restart)")

    pages = [
        ("0 — Overview",           NAVY,  "KPI cards from model_registry.json  ·  Section 10 leaderboard notes  ·  live-verified: disk edit of val_f1_macro reflected on next nav with no restart"),
        ("1 — EDA Gallery",        TEAL,  "14 Section-2 charts (chart_exists() + loading='lazy' placeholders for missing files)"),
        ("2 — Leaderboard",        NAVY,  "Live MLflow leaderboard via dcc.Interval + mlflow_client.get_leaderboard()  ·  'Live (MLflow)' / 'Static fallback' badge"),
        ("3 — Live Predictions",   TEAL,  "All 4 endpoints called in parallel (ThreadPoolExecutor, max_workers=4)  ·  dbc.Alert warning banner when model_status == 'below_quality_bar'"),
        ("4 — Drift Monitoring",   NAVY,  "3-state UI: Prometheus unreachable → static fallback  ·  -1 → 'no check yet'  ·  >0 → live PSI bar chart vs DRIFT_PSI_THRESHOLD"),
        ("5 — Clustering & SHAP",  TEAL,  "K-Means/PCA/t-SNE charts  ·  SHAP beeswarm/waterfall  ·  DistilBERT attention heatmaps  ·  fairness breakdown by Gender/Age/Channel"),
    ]
    y = Inches(1.5)
    for name, color, desc in pages:
        add_rect(s, Inches(0.3), y, Inches(2.9), Inches(0.85), color)
        add_rect(s, Inches(3.3), y, Inches(9.6), Inches(0.85), WHITE)
        txb(s, name, Inches(0.38), y + Inches(0.18),
            Inches(2.7), Inches(0.55), size=13, bold=True, color=WHITE)
        txb(s, desc, Inches(3.38), y + Inches(0.15),
            Inches(9.4), Inches(0.65), size=13)
        y += Inches(0.95)

    bullet_box(s, [
        "Tests: tests/test_dash_app.py — 28 tests covering page registry, all 6 layout() renders with zero data files (3-tier fallback), /charts/ traversal-safety, api_client/mlflow_client mocked success+failure",
        "Docker: docker/dash/Dockerfile — gunicorn WSGI, CMD gunicorn dash_app.app:server --workers 2 --threads 4 --bind 0.0.0.0:8050 --timeout 30  ·  CSIP_API_KEY only (no ADMIN_API_KEY in Dash container)",
    ], Inches(0.3), Inches(7.15), Inches(12.7), Inches(0.65), size=12, color=GRAY)


def slide_findings(prs):
    s = blank_slide(prs)
    bg(s)
    header_bar(s, "Key Findings & Business Insights")

    add_rect(s, Inches(0.4), Inches(1.5), Inches(5.8), Inches(5.6), WHITE)
    txb(s, "ML Results", Inches(0.5), Inches(1.55),
        Inches(5.6), Inches(0.4), size=15, bold=True, color=NAVY)

    rows = [
        ("Task", "Best Model", "Metric", True),
        ("Ticket Type (5-class)", "LightGBM", "F1 = 0.1997 ≈ 0.20 noise floor", False),
        ("Ticket Type (5-class)", "DistilBERT", "F1 = 0.1954 ≈ same noise floor", False),
        ("Ticket Priority (4-class)", "XGBoost ✅", "F1 = 0.2625  (real signal)", False),
        ("Resolution Time", "All models", "RMSE ≈ 7.09 hrs  R² ≈ −0.02", False),
        ("Resolution Time", "Dummy mean", "RMSE = 7.06 hrs  (marginal winner)", False),
    ]
    y = Inches(1.95)
    for row in rows:
        hdr = row[3]
        fc  = TEAL if hdr else (LIGHT if rows.index(row) % 2 == 0 else WHITE)
        tc  = WHITE if hdr else DARK
        for cell, cx, cw in zip(row[:3],
                                 [Inches(0.5), Inches(2.5), Inches(3.8)],
                                 [Inches(1.95), Inches(1.25), Inches(2.35)]):
            add_rect(s, cx, y, cw, Inches(0.5), fc)
            txb(s, cell, cx + Inches(0.05), y + Inches(0.08),
                cw - Inches(0.08), Inches(0.4), size=12, bold=hdr, color=tc)
        y += Inches(0.5)

    add_rect(s, Inches(6.5), Inches(1.5), Inches(6.4), Inches(5.6), WHITE)
    txb(s, "Business Insights", Inches(6.6), Inches(1.55),
        Inches(6.2), Inches(0.4), size=15, bold=True, color=NAVY)
    bullet_box(s, [
        "1. Automated priority prediction delivers real value",
        "   XGBoost F1=0.2625 vs 0.25 baseline — SHAP confirms",
        "   days_since_purchase as a sensible, explainable driver",
        "   (newer purchases → higher urgency)",
        "",
        "2. 80% SLA breach rate for Critical tickets",
        "   Strong operational case for automated triage even",
        "   before model accuracy improves further",
        "",
        "3. Ticket Type labels carry no learnable signal",
        "   4 independent methods confirm — likely synthetic",
        "   assignment in the source dataset",
        "   → API signals this honestly (model_status field)",
        "   → Not hidden or silently worked around",
        "",
        "4. Engineering rigor caught real bugs at every layer",
        "   Preprocessor.pkl wrapper, host port conflict,",
        "   scheduler reconnect failure, stale Docker image —",
        "   all found only by actually running the system",
    ], Inches(6.6), Inches(1.95), Inches(6.2), Inches(5.0), size=13)


def slide_conclusion(prs):
    s = blank_slide(prs)
    bg(s, NAVY)

    add_rect(s, 0, Inches(1.0), W, Inches(5.5), RGBColor(0x0D, 0x25, 0x40))

    txb(s, "What We Built", Inches(0.6), Inches(1.15), Inches(12), Inches(0.55),
        size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    achieved = [
        "✅ End-to-end ML system covering all 3 business tasks",
        "✅ DistilBERT fine-tuning on Colab T4  ·  SHAP + attention explainability",
        "✅ Optuna HPO (6 studies)  ·  MLflow tracking (47+ runs, model registry)",
        "✅ FastAPI serving layer with auth, hot-reload, Prometheus metrics",
        "✅ Prometheus + Grafana + Alertmanager + Evidently drift detection + RUNBOOK",
        "✅ 4 Airflow DAGs with PSI-gated per-task retraining (26 tasks)",
        "✅ 10-container Docker Compose full-stack deployment",
        "✅ 6-page Plotly Dash dashboard  ·  55 pytest tests  ·  GitHub Actions CI",
        "✅ Fairness analysis  ·  Dummy baselines  ·  PEP8/flake8 compliance",
    ]
    bullet_box(s, achieved, Inches(1.5), Inches(1.8), Inches(10.5), Inches(4.0),
               size=15, color=WHITE)

    txb(s, "GitHub:  github.com/PriyaMonisha/AI-Powered-Customer-Support-Intelligence-Platform",
        Inches(0.6), Inches(6.4), Inches(12), Inches(0.5),
        size=14, color=RGBColor(0x85, 0xC1, 0xE9), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    slide_title(prs)
    slide_agenda(prs)
    slide_problem(prs)
    slide_dataset(prs)
    slide_architecture(prs)
    slide_eda(prs)
    slide_features(prs)
    slide_models(prs)
    slide_distilbert(prs)
    slide_regression_clustering(prs)
    slide_shap(prs)
    slide_mlflow(prs)
    slide_fastapi(prs)
    slide_monitoring(prs)
    slide_airflow(prs)
    slide_docker(prs)
    slide_dashboard(prs)
    slide_findings(prs)
    slide_conclusion(prs)

    out = os.path.join(os.path.dirname(os.path.dirname(__file__)),
                       "CSIP_Presentation.pptx")
    prs.save(out)
    print(f"Saved {len(prs.slides)} slides -> {out}")


if __name__ == "__main__":
    main()
